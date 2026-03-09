from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("output/ppt")
OUT_FILE = OUT_DIR / "tenderwriter_presentazione_manageriale_tecnica.pptx"
SUBTITLE = "Analisi basata sul codebase TenderWriter - snapshot locale del repository"
FOOTER = "Fonte: repository, README, docker-compose, backend/app, frontend/src"

COL = {
    "navy": RGBColor(15, 23, 42),
    "slate": RGBColor(71, 85, 105),
    "muted": RGBColor(100, 116, 139),
    "light": RGBColor(248, 250, 252),
    "line": RGBColor(203, 213, 225),
    "white": RGBColor(255, 255, 255),
    "blue": RGBColor(37, 99, 235),
    "blue_soft": RGBColor(219, 234, 254),
    "green": RGBColor(22, 163, 74),
    "green_soft": RGBColor(220, 252, 231),
    "teal": RGBColor(13, 148, 136),
    "teal_soft": RGBColor(204, 251, 241),
    "amber": RGBColor(217, 119, 6),
    "amber_soft": RGBColor(254, 243, 199),
    "rose": RGBColor(190, 24, 93),
    "rose_soft": RGBColor(251, 207, 232),
    "violet": RGBColor(109, 40, 217),
    "violet_soft": RGBColor(237, 233, 254),
}


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text="", size=20, bold=False, color=None, align=PP_ALIGN.LEFT,
            fill=None, line=None, shape=None, valign=MSO_ANCHOR.TOP):
    shp = slide.shapes.add_shape(shape, x, y, w, h) if shape else slide.shapes.add_textbox(x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or COL["navy"]
    return shp


def title(slide, text, subtitle=None, dark=False):
    textbox(slide, Inches(0.6), Inches(0.35), Inches(11.4), Inches(0.6), text, 26, True, COL["white"] if dark else COL["navy"])
    if subtitle:
        textbox(slide, Inches(0.6), Inches(0.95), Inches(11.2), Inches(0.35), subtitle, 10.5, False, COL["line"] if dark else COL["slate"])
    if not dark:
        textbox(slide, Inches(0.6), Inches(1.42), Inches(1.2), Inches(0.08), "", fill=COL["blue"], shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)


def footer(slide, n):
    textbox(slide, Inches(0.6), Inches(6.95), Inches(10.1), Inches(0.22), FOOTER, 8, False, COL["muted"])
    textbox(slide, Inches(12.0), Inches(6.95), Inches(0.5), Inches(0.22), str(n), 8, False, COL["muted"], PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, size=16, color=None):
    shp = slide.shapes.add_textbox(x, y, w, h)
    shp.fill.background()
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(7)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COL["navy"]
    return shp


def card(slide, x, y, w, h, head, body, fill, line):
    shp = textbox(slide, x, y, w, h, "", fill=fill, line=line, shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE)
    tf = shp.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = head
    r1.font.name = "Aptos"
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = COL["navy"]
    p2 = tf.add_paragraph()
    p2.space_before = Pt(5)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Aptos"
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = COL["slate"]
    return shp


def metric(slide, x, y, w, h, num, label, accent):
    shp = textbox(slide, x, y, w, h, "", fill=COL["white"], line=COL["line"], shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE)
    textbox(slide, x, y, Inches(0.08), h, "", fill=accent, shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    tf = shp.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = str(num)
    r1.font.name = "Aptos"
    r1.font.size = Pt(23)
    r1.font.bold = True
    r1.font.color.rgb = COL["navy"]
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Aptos"
    r2.font.size = Pt(10)
    r2.font.color.rgb = COL["slate"]
    return shp


def conn(slide, x1, y1, x2, y2, color):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    return line


def section(prs, text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["navy"])
    textbox(slide, Inches(0.75), Inches(0.95), Inches(2.2), Inches(0.4), "SEZIONE", 11, True, COL["white"], fill=COL["blue"], shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE)
    textbox(slide, Inches(0.75), Inches(2.0), Inches(8.5), Inches(0.9), text, 28, True, COL["white"])
    textbox(slide, Inches(0.75), Inches(3.0), Inches(8.5), Inches(0.45), subtitle_text, 14, False, COL["line"])
    textbox(slide, Inches(0.75), Inches(6.7), Inches(9.0), Inches(0.2), SUBTITLE, 9, False, COL["line"])
    return slide


def flow(slide, x, y, w, h, num, head, body, fill, line):
    bubble = textbox(slide, x, y, Inches(0.38), Inches(0.38), str(num), 12, True, COL["white"], PP_ALIGN.CENTER, line, line, MSO_AUTO_SHAPE_TYPE.OVAL, MSO_ANCHOR.MIDDLE)
    card_box = card(slide, x, y + Inches(0.48), w, h, head, body, fill, line)
    return bubble, card_box


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    page = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["navy"])
    textbox(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.18), "", fill=RGBColor(10, 15, 31), shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    textbox(slide, Inches(0.75), Inches(1.48), Inches(2.5), Inches(0.45), "MANAGERIAL + TECHNICAL DECK", 11, True, COL["white"], fill=COL["teal"], shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE)
    textbox(slide, Inches(0.75), Inches(2.15), Inches(8.6), Inches(1.15), "TenderWriter\nPiattaforma AI locale per la gestione e la scrittura di proposte di gara", 26, True, COL["white"])
    textbox(slide, Inches(0.77), Inches(3.55), Inches(7.6), Inches(0.55), "Presentazione professionale basata sul codebase: architettura, funzionalita' realizzate, processo di sviluppo e roadmap.", 14, False, COL["line"])
    textbox(slide, Inches(0.77), Inches(6.65), Inches(6.0), Inches(0.25), "Snapshot analizzato: 8 marzo 2026", 9.5, False, COL["line"])
    textbox(slide, Inches(11.1), Inches(6.65), Inches(1.5), Inches(0.25), "Codex generated", 9.5, False, COL["line"], PP_ALIGN.RIGHT)
    card(slide, Inches(9.0), Inches(1.95), Inches(3.35), Inches(1.0), "Tesi centrale", "TenderWriter e' gia' impostato come piattaforma local-first capace di coprire la catena chiave di tender management, content generation e collaborative editing.", COL["white"], COL["blue"])
    metric(slide, Inches(9.0), Inches(3.2), Inches(1.55), Inches(0.92), 16, "servizi orchestrati", COL["blue"])
    metric(slide, Inches(10.72), Inches(3.2), Inches(1.55), Inches(0.92), 9, "famiglie API", COL["teal"])
    metric(slide, Inches(9.0), Inches(4.28), Inches(1.55), Inches(0.92), 10, "aree UI", COL["amber"])
    metric(slide, Inches(10.72), Inches(4.28), Inches(1.55), Inches(0.92), 11, "entita' dati", COL["rose"])
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Executive summary", SUBTITLE)
    for args in [
        (0.7, 1.8, 3.9, 1.35, "Posizionamento", "Suite open-source per team che preparano offerte e vogliono usare AI e retrieval su infrastruttura locale, non dipendente dal cloud.", COL["white"], COL["line"]),
        (4.75, 1.8, 3.9, 1.35, "Capacita' reali", "Autenticazione con OTP, CRUD tender/proposte, content library, ricerca AI, editing OnlyOffice, task asincroni e monitoraggio runtime.", COL["white"], COL["line"]),
        (8.8, 1.8, 3.9, 1.35, "Stato attuale", "Prodotto in sviluppo avanzato: base piattaforma solida, ampia copertura funzionale, margini di maturazione su roadmap e industrializzazione.", COL["white"], COL["line"]),
    ]:
        card(slide, Inches(args[0]), Inches(args[1]), Inches(args[2]), Inches(args[3]), *args[4:])
    bullets(slide, Inches(0.78), Inches(3.55), Inches(12.0), Inches(2.55), [
        "La codebase mostra una separazione architetturale chiara: frontend React, backend FastAPI, servizi dati specializzati, task async e due runtime LLM distinti.",
        "Il valore business e' nell'accorciare il ciclo di lavoro delle gare: ingestione documenti, estrazione requisiti, generazione sezioni, collaborazione e controllo dell'infrastruttura.",
        "L'impostazione local-first e' coerente con scenari sensibili su dati e compliance e riduce il lock-in verso servizi cloud.",
        "La struttura del progetto consente un'evoluzione incrementale per workstream: security, RAG, editing, observability e governance amministrativa.",
    ], 17)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Problema, utenti target e proposta di valore", "Lettura manageriale del prodotto")
    for args in [
        (0.7, 1.7, 3.8, 1.45, "Problema", "La preparazione di una gara richiede coordinamento di contenuti, requisiti, dati storici, team e versioni, spesso con strumenti frammentati.", COL["blue_soft"], COL["blue"]),
        (4.76, 1.7, 3.8, 1.45, "Utenti prioritari", "Bid manager, proposal manager, redattori tecnici, admin di piattaforma e figure IT che devono governare ambiente e servizi.", COL["teal_soft"], COL["teal"]),
        (8.82, 1.7, 3.8, 1.45, "Valore", "Riduzione del tempo di preparazione, maggior riuso dei contenuti, tracciabilita' dei requisiti e controllo dell'infrastruttura AI.", COL["amber_soft"], COL["amber"]),
    ]:
        card(slide, Inches(args[0]), Inches(args[1]), Inches(args[2]), Inches(args[3]), *args[4:])
    for args in [
        (0.75, 3.55, 2.6, 1.0, "local-first", "modello operativo", COL["teal"]),
        (3.6, 3.55, 2.6, 1.0, "HybridRAG", "motore AI", COL["blue"]),
        (6.45, 3.55, 2.6, 1.0, "OnlyOffice", "collaboration layer", COL["rose"]),
        (9.3, 3.55, 2.6, 1.0, "Celery + Redis", "execution model", COL["amber"]),
    ]:
        metric(slide, Inches(args[0]), Inches(args[1]), Inches(args[2]), Inches(args[3]), args[4], args[5], args[6])
    bullets(slide, Inches(0.8), Inches(5.05), Inches(12.0), Inches(1.5), [
        "La proposta e' credibile per organizzazioni che richiedono controllo su dati, modelli e stack infrastrutturale.",
        "La codebase supporta sia il punto di vista operativo di business sia il punto di vista amministrativo/tecnico con pagine dedicate.",
        "L'integrazione di OpenCode mostra un'estensione naturale verso agentic coding e manutenzione locale della piattaforma.",
    ], 16)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Funzionalita' realizzate", "Copertura funzionale osservabile nel repository")
    func_cards = [
        ("Identity & access", "Registrazione utente, verifica OTP, login JWT, ruoli admin/editor, RBAC per tender.", COL["blue_soft"], COL["blue"]),
        ("Tender lifecycle", "Creazione tender, import documenti, stati di avanzamento, requisiti e permessi granulati.", COL["teal_soft"], COL["teal"]),
        ("Proposal workspace", "Creazione proposta, sezioni di default, aggiornamento contenuti e workflow di stato.", COL["green_soft"], COL["green"]),
        ("AI assistance", "Query RAG, generazione sezione, compliance check, analisi requisiti e cronologia ricerche.", COL["violet_soft"], COL["violet"]),
        ("Content reuse", "Libreria blocchi contenuto, tagging, categorie, quality rating e riuso operativo.", COL["amber_soft"], COL["amber"]),
        ("Operations & admin", "System monitor, stats container, logs, gestione utenti e permessi, timeout Nginx.", COL["rose_soft"], COL["rose"]),
    ]
    pos = [(0.7, 1.75), (4.5, 1.75), (8.3, 1.75), (0.7, 4.1), (4.5, 4.1), (8.3, 4.1)]
    for (head, body, fill, line), (x, y) in zip(func_cards, pos):
        card(slide, Inches(x), Inches(y), Inches(3.2), Inches(1.85), head, body, fill, line)
    footer(slide, page)
    page += 1

    section(prs, "Approfondimento tecnico", "Architettura, flussi, moduli e processo di engineering ricostruiti dalla codebase.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Architettura di piattaforma", "Vista logica a layer")
    cols = ["Experience layer", "Application layer", "AI & workflow layer", "Data layer", "Ops layer"]
    xs = [0.75, 3.2, 5.65, 8.1, 10.55]
    for x, txt in zip(xs, cols):
        textbox(slide, Inches(x), Inches(1.75), Inches(2.15), Inches(0.3), txt, 12, True, COL["slate"])
    arch_cards = [
        ("Frontend React", "Dashboard, proposals, library, AI search, tasks, monitor e settings.", COL["green_soft"], COL["green"]),
        ("FastAPI backend", "Router modulari: auth, tenders, proposals, content, RAG, admin, system, tasks, OnlyOffice.", COL["blue_soft"], COL["blue"]),
        ("HybridRAG + Celery", "Retrieval multi-strategia, generazione LLM, task asincroni e schedule di cleanup.", COL["violet_soft"], COL["violet"]),
        ("Postgres / Qdrant / Neo4j / Redis / MinIO", "Persistenza transazionale, vettoriale, grafo, broker e object storage.", COL["amber_soft"], COL["amber"]),
        ("OnlyOffice / Mailpit / Redis Insight / OpenCode", "Servizi satellite per editing, mail dev, inspection e coding agent.", COL["rose_soft"], COL["rose"]),
    ]
    blocks = []
    for x, (head, body, fill, line) in zip(xs, arch_cards):
        blocks.append(card(slide, Inches(x), Inches(2.1), Inches(2.0), Inches(1.1), head, body, fill, line))
    for i, color in enumerate([COL["slate"], COL["blue"], COL["amber"], COL["rose"]]):
        conn(slide, blocks[i].left + blocks[i].width, blocks[i].top + Inches(0.55), blocks[i + 1].left, blocks[i + 1].top + Inches(0.55), color)
    bullets(slide, Inches(0.8), Inches(3.75), Inches(12.0), Inches(2.1), [
        "Il deployment e' organizzato in Docker Compose con 16 servizi e dipendenze esplicite tra backend, data store, worker e LLM.",
        "La piattaforma distingue il motore AI di prodotto (`llama-tender`) da quello dedicato all'agente di coding (`llama-opencode`).",
        "La scelta di servizi dati specializzati rende l'architettura estensibile: dominio, semantica, relazioni, broker e object storage.",
        "Il backend agisce anche come control plane operativo, non solo come API applicativa.",
    ], 16)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Flusso end-to-end del prodotto", "Dal documento di gara alla proposta collaborativa")
    flow_items = [
        ("Tender import", "Upload del documento gara e memorizzazione in object storage.", COL["blue_soft"], COL["blue"]),
        ("Requirement extraction", "Parsing, chunking e analisi requisiti tramite pipeline documentale.", COL["teal_soft"], COL["teal"]),
        ("Proposal setup", "Creazione proposta e struttura sezioni iniziali dal backend.", COL["green_soft"], COL["green"]),
        ("AI assisted writing", "Ricerca RAG, generazione sezioni e compliance check con LLM locale.", COL["violet_soft"], COL["violet"]),
        ("Collaborative editing", "Editing con OnlyOffice, callback di salvataggio e re-indicizzazione.", COL["rose_soft"], COL["rose"]),
        ("Async export & ops", "Task manager, export PDF, monitoraggio servizi e cleanup schedulati.", COL["amber_soft"], COL["amber"]),
    ]
    lefts = [0.55, 2.7, 4.85, 7.0, 9.15, 11.3]
    cards = []
    for idx, ((head, body, fill, line), left) in enumerate(zip(flow_items, lefts), start=1):
        cards.append(flow(slide, Inches(left), Inches(2.1), Inches(1.93), Inches(2.05), idx, head, body, fill, line)[1])
    for i in range(len(cards) - 1):
        conn(slide, cards[i].left + cards[i].width, cards[i].top + Inches(0.7), cards[i + 1].left, cards[i + 1].top + Inches(0.7), COL["slate"])
    bullets(slide, Inches(0.8), Inches(5.25), Inches(12.0), Inches(1.25), [
        "Il flusso e' visibile nelle API, nei task, nelle pagine frontend e nelle integrazioni infrastrutturali.",
        "La pipeline combina lavoro sincrono per l'utente e lavorazioni asincrone per attivita' costose o batch.",
    ], 16)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Deep dive: motore HybridRAG", "Elemento differenziante della piattaforma")
    rag_cards = [
        ("Dense retrieval", "Ricerca semantica in Qdrant su embedding dei chunk documentali.", COL["blue_soft"], COL["blue"]),
        ("Sparse retrieval", "BM25 in memoria per keyword, sigle, codici e termini tecnici.", COL["teal_soft"], COL["teal"]),
        ("Graph retrieval", "Neo4j per relazioni tra progetti, team member, certificazioni e requirement.", COL["amber_soft"], COL["amber"]),
        ("Fusion & reranking", "Unione dei risultati via reciprocal rank fusion e filtro finale piu' preciso.", COL["violet_soft"], COL["violet"]),
        ("Generation", "Template di prompt e generazione con llama.cpp server.", COL["rose_soft"], COL["rose"]),
    ]
    for i, (head, body, fill, line) in enumerate(rag_cards):
        card(slide, Inches(0.7 + i * 2.3), Inches(1.75), Inches(2.1), Inches(1.3), head, body, fill, line)
    chart_data = CategoryChartData()
    chart_data.categories = ["Dense", "Sparse", "Graph", "Fusion", "Generation"]
    chart_data.add_series("Pipeline", (1, 1, 1, 1, 1))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(3.45), Inches(4.2), Inches(2.2), chart_data).chart
    chart.has_legend = False
    chart.value_axis.visible = False
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = COL["blue"]
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Pipeline HybridRAG implementata"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    bullets(slide, Inches(5.35), Inches(3.35), Inches(7.1), Inches(2.4), [
        "L'engine supporta modalita' diverse: `search`, `qa`, `write_section`, `exec_summary`, `analyze_reqs`, `compliance`.",
        "La query puo' essere anche streamata e viene salvata nella `search_history` dell'utente autenticato.",
        "La piattaforma usa il motore RAG sia per la ricerca sia per la scrittura guidata delle sezioni di proposta.",
        "Il layer di ingestion alimenta vettori, BM25 e knowledge graph, aumentando resilienza e quality of answer.",
    ], 16)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Moduli applicativi e responsabilita'", "Vista tecnica per stakeholder IT e architetturali")
    mod_cards = [
        ("Frontend workspace", "Routing React con aree dedicate a dashboard, proposals, content library, AI search, tasks, monitor e settings.", COL["green_soft"], COL["green"]),
        ("Backend domain API", "Router separati per auth, tenders, proposals, content blocks, admin, system, tasks e OnlyOffice.", COL["blue_soft"], COL["blue"]),
        ("Domain model", "Entita' per utenti, tender, requirements, proposal, sections, content block, document, chunk e permission.", COL["amber_soft"], COL["amber"]),
        ("Document collaboration", "OnlyOffice genera config, token JWT, callback di save, force-save e re-indicizzazione nel RAG.", COL["rose_soft"], COL["rose"]),
        ("Async processing", "Celery + Redis per indexing, generation, export PDF, health check e cleanup schedulati.", COL["violet_soft"], COL["violet"]),
        ("Admin & observability", "Docker stats/logs, lista componenti, controllo runtime e pagina Developments con changelog sintetico.", COL["teal_soft"], COL["teal"]),
    ]
    pos = [(0.75, 1.7), (4.48, 1.7), (8.21, 1.7), (0.75, 4.1), (4.48, 4.1), (8.21, 4.1)]
    for (head, body, fill, line), (x, y) in zip(mod_cards, pos):
        card(slide, Inches(x), Inches(y), Inches(3.1), Inches(1.9), head, body, fill, line)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Processo di sviluppo osservabile nel codebase", "Come il prodotto viene costruito, evoluto e governato")
    dev_steps = [
        ("Local setup", "Avvio completo via Docker Compose con servizi, dati, mail dev e runtime AI su macchina locale.", COL["blue_soft"], COL["blue"]),
        ("Foundation", "Sicurezza, config validation, sessioni, ruoli, persistence e monitoraggio introdotti presto.", COL["teal_soft"], COL["teal"]),
        ("Feature increments", "Workstream separati per Celery, MinIO, Task Manager, Section Management, OTP hardening e admin tools.", COL["green_soft"], COL["green"]),
        ("Verification", "Script di verifica, dipendenze dev, linter/test e pagine di sviluppo supportano debug e quality control.", COL["amber_soft"], COL["amber"]),
        ("Operational feedback", "System Monitor, Task Manager, logs e metriche forniscono un ciclo di feedback breve durante lo sviluppo.", COL["rose_soft"], COL["rose"]),
    ]
    xs = [0.75, 3.2, 5.65, 8.1, 10.55]
    step_cards = []
    for i, ((head, body, fill, line), x) in enumerate(zip(dev_steps, xs), start=1):
        step_cards.append(flow(slide, Inches(x), Inches(1.75), Inches(2.1), Inches(1.55), i, head, body, fill, line)[1])
    for i in range(len(step_cards) - 1):
        conn(slide, step_cards[i].left + step_cards[i].width, step_cards[i].top + Inches(0.53), step_cards[i + 1].left, step_cards[i + 1].top + Inches(0.53), COL["slate"])
    bullets(slide, Inches(0.8), Inches(4.4), Inches(12.0), Inches(1.9), [
        "2026-03-06: completati Celery integration, MinIO storage, Task Manager UI, Components dashboard, rate limiting, OTP hardening e secure env validation.",
        "2026-03-07: pagina Developments segnala Section Management in progress e conferma un approccio incrementale per feature branch/workstream.",
        "README e compose mostrano un modello pragmatico: debug guidato da log, rebuild mirati, servizi mock e roadmap esplicita.",
    ], 15.5)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Qualita', sicurezza e operabilita'", "Valutazione tecnica del livello di maturita'")
    quality_cards = [
        ("Security baseline", "OTP con limiti tentativi, JWT, password hashing, rate limiting SlowAPI, validazione secret e ruoli amministrativi.", COL["blue_soft"], COL["blue"]),
        ("Operational control", "Health endpoint, RAG health, log container, stats CPU/RAM, hot update timeout Nginx e cleanup schedulati.", COL["teal_soft"], COL["teal"]),
        ("Developer tooling", "Ruff, mypy, pytest, pytest-asyncio, Vitest, verify_imports, e2e_verify e stack dev riproducibile.", COL["amber_soft"], COL["amber"]),
        ("Risk note", "La maturita' operativa e' promettente ma non ancora industrializzata.", COL["rose_soft"], COL["rose"]),
    ]
    dims = [(0.75, 1.72, 3.05), (4.12, 1.72, 3.05), (7.49, 1.72, 3.05), (10.86, 1.72, 1.72)]
    for (head, body, fill, line), (x, y, w) in zip(quality_cards, dims):
        card(slide, Inches(x), Inches(y), Inches(w), Inches(1.7), head, body, fill, line)
    bullets(slide, Inches(0.85), Inches(4.0), Inches(12.0), Inches(2.25), [
        "Punto di forza: l'architettura incorpora gia' sicurezza applicativa, observability e gestione operativa, non solo feature AI.",
        "Punto di attenzione: alcune parti mostrano integrazioni ancora da consolidare, soprattutto nella piena industrializzazione dei task e nella chiusura della roadmap.",
        "Dal punto di vista manageriale, la base e' adatta a un programma di hardening mirato, non a una ripartenza da zero.",
    ], 16)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["light"])
    title(slide, "Stato attuale, gap e roadmap", "Lettura di priorita' per decision makers")
    roadmap_cards = [
        ("Gia' realizzato", "Piattaforma end-to-end locale con autenticazione, RAG, gestione tender/proposte, content library, task async, OnlyOffice e monitoraggio.", COL["green_soft"], COL["green"]),
        ("Gap di prodotto", "Integrazione completa della ricerca con cronologia utente, export professionale PDF/Docx e compliance matrix piu' evoluta restano in roadmap.", COL["amber_soft"], COL["amber"]),
        ("Gap di industrializzazione", "Servono test coverage piu' strutturata, rifinitura UX, hardening deployment e una chiara readiness production.", COL["rose_soft"], COL["rose"]),
    ]
    for (head, body, fill, line), x in zip(roadmap_cards, [0.75, 4.75, 8.75]):
        card(slide, Inches(x), Inches(1.7), Inches(3.85), Inches(1.8), head, body, fill, line)
    chart_data = CategoryChartData()
    chart_data.categories = ["Core platform", "AI retrieval", "Collab editing", "Ops/admin", "Production hardening"]
    chart_data.add_series("Maturita' relativa", (85, 80, 78, 76, 55))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.85), Inches(4.05), Inches(5.5), Inches(2.0), chart_data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = COL["blue"]
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Lettura di maturita' ricostruita dal codebase"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    bullets(slide, Inches(6.65), Inches(4.1), Inches(5.6), Inches(1.95), [
        "Decisione suggerita: investire in un ciclo breve di product hardening, non in un redesign completo.",
        "Priorita' 1: consolidare export e compliance workflow.",
        "Priorita' 2: migliorare readiness operativa e test.",
        "Priorita' 3: raffinare experience e analytics utente.",
    ], 15.5)
    footer(slide, page)
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, COL["navy"])
    title(slide, "Conclusioni", "Sintesi finale per sponsor, management e team tecnico", True)
    bullets(slide, Inches(0.9), Inches(1.8), Inches(8.0), Inches(3.2), [
        "TenderWriter non e' solo un prototipo AI: il repository mostra una piattaforma strutturata, con dominio, retrieval, editing collaborativo e operations.",
        "Il design local-first e la modularita' dei servizi sono coerenti con contesti enterprise sensibili e con esigenze di controllo architetturale.",
        "La prossima leva di valore e' portare la base esistente a uno standard di affidabilita', testabilita' e presentazione commerciale superiore.",
    ], 19, COL["white"])
    card(slide, Inches(9.0), Inches(1.95), Inches(3.1), Inches(1.1), "Messaggio chiave", "La base tecnica e' sufficientemente solida da giustificare un investimento di consolidamento e go-to-market interno.", COL["white"], COL["blue"])
    card(slide, Inches(9.0), Inches(3.35), Inches(3.1), Inches(1.1), "Uso del deck", "Adatto a steering committee, sponsor IT, partner tecnici e sessioni di allineamento interno.", COL["white"], COL["teal"])
    textbox(slide, Inches(0.9), Inches(6.7), Inches(5.0), Inches(0.25), SUBTITLE, 9, False, COL["line"])

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))
    return OUT_FILE


if __name__ == "__main__":
    print(build_presentation())
